"""
Générateur de présentation PowerPoint pour le projet Jarvim / MCP
Charte graphique : eDEX-UI Futuristic Terminal Aesthetic
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── Palette de couleurs (d'après index.css) ───────────────────────────────────
BG       = RGBColor(0x05, 0x0D, 0x0C)   # oklch(0.070 0.020 185) - noir teal profond
FG       = RGBColor(0xD8, 0xF5, 0xF0)   # oklch(0.910 0.055 175) - blanc cyan clair
CYAN     = RGBColor(0x00, 0xE8, 0xDA)   # oklch(0.880 0.195 178) - cyan primaire
YELLOW   = RGBColor(0xC5, 0xAF, 0x00)   # oklch(0.720 0.185 60)  - jaune accent
DARK2    = RGBColor(0x0D, 0x1A, 0x1F)   # oklch(0.130 0.030 200) - bleu-gris sombre
RED      = RGBColor(0xE5, 0x40, 0x28)   # oklch(0.600 0.230 25)  - rouge destructif
MUTED    = RGBColor(0x2A, 0x4A, 0x45)   # version atténuée du cyan

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # layout vide


# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(18), bold=False, color=FG,
             align=PP_ALIGN.LEFT, font_name="Courier New"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bg(slide):
    """Fond sombre + ligne de scan subtile."""
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=BG)


def add_corner_brackets(slide):
    """Ajoute des coins décoratifs style terminal."""
    bracket_color = CYAN
    thickness = Pt(2)
    size = Inches(0.35)
    margin = Inches(0.15)

    corners = [
        # top-left
        [(margin, margin, size, Inches(0.03)),
         (margin, margin, Inches(0.03), size)],
        # top-right
        [(SLIDE_W - margin - size, margin, size, Inches(0.03)),
         (SLIDE_W - margin - Inches(0.03), margin, Inches(0.03), size)],
        # bottom-left
        [(margin, SLIDE_H - margin - Inches(0.03), size, Inches(0.03)),
         (margin, SLIDE_H - margin - size, Inches(0.03), size)],
        # bottom-right
        [(SLIDE_W - margin - size, SLIDE_H - margin - Inches(0.03), size, Inches(0.03)),
         (SLIDE_W - margin - Inches(0.03), SLIDE_H - margin - size, Inches(0.03), size)],
    ]
    for corner in corners:
        for (x, y, w, h) in corner:
            add_rect(slide, x, y, w, h, fill_color=bracket_color)


def add_header_bar(slide, title_text, subtitle_text=None):
    """Barre d'en-tête avec titre style terminal."""
    # Barre de fond
    add_rect(slide, Inches(0.5), Inches(0.15),
             SLIDE_W - Inches(1), Inches(0.75),
             fill_color=DARK2)
    # Ligne cyan sous la barre
    add_rect(slide, Inches(0.5), Inches(0.90),
             SLIDE_W - Inches(1), Inches(0.03),
             fill_color=CYAN)

    # Préfixe terminal
    prefix = add_text(slide, "> ", Inches(0.55), Inches(0.18),
                      Inches(0.4), Inches(0.65),
                      font_size=Pt(22), bold=True, color=CYAN)

    # Titre
    add_text(slide, title_text, Inches(0.95), Inches(0.18),
             Inches(9), Inches(0.65),
             font_size=Pt(22), bold=True, color=FG)

    if subtitle_text:
        add_text(slide, subtitle_text, Inches(0.95), Inches(0.18),
                 SLIDE_W - Inches(1.5), Inches(0.65),
                 font_size=Pt(14), color=YELLOW, align=PP_ALIGN.RIGHT)


def add_footer(slide, slide_number, total=9):
    """Pied de page style terminal."""
    add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.55),
             SLIDE_W - Inches(1), Inches(0.03),
             fill_color=MUTED)
    add_text(slide, f"JARVIM MCP // 2025-2026", Inches(0.55), SLIDE_H - Inches(0.52),
             Inches(6), Inches(0.4), font_size=Pt(10), color=MUTED)
    add_text(slide, f"[{slide_number:02d}/{total:02d}]", SLIDE_W - Inches(1.8), SLIDE_H - Inches(0.52),
             Inches(1.5), Inches(0.4), font_size=Pt(10), color=CYAN, align=PP_ALIGN.RIGHT)


def add_section_label(slide, label, x, y):
    """Label de section style terminal."""
    add_text(slide, f"# {label}", x, y, Inches(4), Inches(0.35),
             font_size=Pt(11), color=YELLOW, bold=True)


def divider_line(slide, y):
    add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1), Inches(0.02), fill_color=MUTED)


def add_tag(slide, text, x, y, w=Inches(2), h=Inches(0.42)):
    """Petit badge style tag."""
    add_rect(slide, x, y, w, h, fill_color=DARK2, line_color=CYAN, line_width=Pt(1))
    add_text(slide, text, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), h - Inches(0.04),
             font_size=Pt(12), color=CYAN, align=PP_ALIGN.CENTER)


def add_feature_card(slide, icon, title, desc, x, y, w=Inches(3.7), h=Inches(1.5)):
    """Carte de fonctionnalité."""
    add_rect(slide, x, y, w, h, fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
    # Ligne cyan gauche
    add_rect(slide, x, y, Inches(0.06), h, fill_color=CYAN)
    add_text(slide, icon + " " + title, x + Inches(0.15), y + Inches(0.1),
             w - Inches(0.3), Inches(0.4), font_size=Pt(13), bold=True, color=CYAN)
    add_text(slide, desc, x + Inches(0.15), y + Inches(0.52),
             w - Inches(0.3), h - Inches(0.62), font_size=Pt(11), color=FG)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INTRODUCTION (Titre)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)

# Ligne déco haut
add_rect(slide, Inches(0.5), Inches(1.8), SLIDE_W - Inches(1), Inches(0.04), fill_color=CYAN)

# Logo texte
add_text(slide, "JARVIM", Inches(0.5), Inches(2.0),
         SLIDE_W - Inches(1), Inches(2.2),
         font_size=Pt(96), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Sous-titre
add_text(slide, "Model Context Protocol  ×  Full-Stack Kanban",
         Inches(0.5), Inches(4.1), SLIDE_W - Inches(1), Inches(0.6),
         font_size=Pt(20), color=FG, align=PP_ALIGN.CENTER)

# Ligne déco bas
add_rect(slide, Inches(0.5), Inches(4.75), SLIDE_W - Inches(1), Inches(0.04), fill_color=CYAN)

# Tags tech
tags = ["NestJS", "React 19", "PostgreSQL", "Prisma", "TailwindCSS", "Claude MCP"]
tag_w = Inches(1.9)
total_w = len(tags) * tag_w + (len(tags) - 1) * Inches(0.15)
start_x = (SLIDE_W - total_w) / 2
for i, tag in enumerate(tags):
    add_tag(slide, tag, start_x + i * (tag_w + Inches(0.15)), Inches(5.1), w=tag_w)

# Footer minimal
add_text(slide, "YNOV — 2025/2026", Inches(0.5), SLIDE_H - Inches(0.65),
         SLIDE_W - Inches(1), Inches(0.4), font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "CONTEXTE", "Slide 02")
add_footer(slide, 2)

# Bloc gauche — contexte scolaire
add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.2),
         fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.04), fill_color=CYAN)

add_section_label(slide, "PROJET SCOLAIRE", Inches(0.65), Inches(1.2))
items = [
    ("Parcours", "3ème année — Ynov Campus"),
    ("Matière", "Module MCP (Model Context Protocol)"),
    ("Objectif", "Intégrer Claude AI dans une app\nfull-stack via le protocole MCP"),
    ("Livrable", "Application web + démo Claude MCP"),
]
y = Inches(1.7)
for label, val in items:
    add_text(slide, f"[{label}]", Inches(0.65), y, Inches(1.5), Inches(0.35),
             font_size=Pt(11), color=YELLOW, bold=True)
    add_text(slide, val, Inches(2.2), y, Inches(3.9), Inches(0.5),
             font_size=Pt(12), color=FG)
    y += Inches(0.8)

# Bloc droite — timeline
add_rect(slide, Inches(6.55), Inches(1.1), Inches(6.2), Inches(5.2),
         fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
add_rect(slide, Inches(6.55), Inches(1.1), Inches(6.2), Inches(0.04), fill_color=CYAN)

add_section_label(slide, "APPROCHE TECHNIQUE", Inches(6.7), Inches(1.2))
points = [
    "Application Kanban collaborative complète",
    "Architecture monorepo (Turbo + npm workspaces)",
    "API REST sécurisée — JWT dual-token",
    "Intégration Claude MCP pour automatiser\nles actions sur les boards",
    "Charte graphique eDEX-UI — terminal futuriste",
]
y = Inches(1.75)
for pt in points:
    add_rect(slide, Inches(6.7), y + Inches(0.12), Inches(0.08), Inches(0.08), fill_color=CYAN)
    add_text(slide, pt, Inches(6.95), y, Inches(5.6), Inches(0.6),
             font_size=Pt(12), color=FG)
    y += Inches(0.82)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DÉFINITION MCP
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "DÉFINITION MCP", "Slide 03")
add_footer(slide, 3)

# Grand titre centré
add_text(slide, "Model Context Protocol",
         Inches(0.5), Inches(1.1), SLIDE_W - Inches(1), Inches(0.7),
         font_size=Pt(30), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

add_text(slide, "Un protocole open-source d'Anthropic permettant aux LLMs\n"
                "de se connecter à des systèmes externes de façon standardisée.",
         Inches(1), Inches(1.85), SLIDE_W - Inches(2), Inches(0.85),
         font_size=Pt(15), color=FG, align=PP_ALIGN.CENTER)

# Ligne de séparation
divider_line(slide, Inches(2.75))

# 3 colonnes — Principe, Composants, Bénéfices
col_w = Inches(3.9)
cols = [
    {
        "title": "PRINCIPE",
        "color": CYAN,
        "items": [
            "Standardise la communication\nentre un LLM et des outils",
            "Remplace les intégrations\npropriétaires ad-hoc",
            "Protocole JSON-RPC 2.0\nsur stdio / SSE / HTTP",
        ]
    },
    {
        "title": "COMPOSANTS",
        "color": YELLOW,
        "items": [
            "MCP Host — le client LLM\n(Claude Desktop, Cursor...)",
            "MCP Server — expose les\ntools, resources, prompts",
            "Transport Layer — stdio\nou SSE pour la comm.",
        ]
    },
    {
        "title": "BÉNÉFICES",
        "color": RGBColor(0x00, 0xE8, 0x70),
        "items": [
            "Interopérabilité totale\nentre agents et services",
            "Sécurité : le LLM ne voit\nque ce que le serveur expose",
            "Extensible : tools, prompts,\nresources, sampling",
        ]
    },
]
for i, col in enumerate(cols):
    cx = Inches(0.5) + i * (col_w + Inches(0.15))
    cy = Inches(2.85)
    add_rect(slide, cx, cy, col_w, Inches(4.0),
             fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
    add_rect(slide, cx, cy, col_w, Inches(0.04), fill_color=col["color"])
    add_text(slide, col["title"], cx + Inches(0.15), cy + Inches(0.1),
             col_w - Inches(0.3), Inches(0.45),
             font_size=Pt(14), bold=True, color=col["color"])
    iy = cy + Inches(0.65)
    for item in col["items"]:
        add_rect(slide, cx + Inches(0.15), iy + Inches(0.1),
                 Inches(0.08), Inches(0.08), fill_color=col["color"])
        add_text(slide, item, cx + Inches(0.35), iy,
                 col_w - Inches(0.5), Inches(0.75), font_size=Pt(11.5), color=FG)
        iy += Inches(0.9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ÉQUIPE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "ÉQUIPE", "Slide 04")
add_footer(slide, 4)

add_text(slide, "Les développeurs du projet",
         Inches(0.5), Inches(1.05), SLIDE_W - Inches(1), Inches(0.5),
         font_size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)

# Carte Enzo
card_w = Inches(5.5)
card_h = Inches(4.2)
# Carte 1
cx = Inches(0.9)
cy = Inches(1.65)
add_rect(slide, cx, cy, card_w, card_h, fill_color=DARK2, line_color=CYAN, line_width=Pt(2))
add_rect(slide, cx, cy, card_w, Inches(0.05), fill_color=CYAN)

# Avatar carré style terminal
add_rect(slide, cx + Inches(1.9), cy + Inches(0.25), Inches(1.7), Inches(1.7),
         fill_color=BG, line_color=CYAN, line_width=Pt(2))
add_text(slide, "ET", cx + Inches(1.9), cy + Inches(0.35), Inches(1.7), Inches(1.5),
         font_size=Pt(44), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

add_text(slide, "ENZO TURPIN", cx, cy + Inches(2.15), card_w, Inches(0.55),
         font_size=Pt(22), bold=True, color=FG, align=PP_ALIGN.CENTER)

infos = [
    ("Rôle", "Lead Developer — Full-Stack"),
    ("Stack", "NestJS · React · PostgreSQL"),
    ("GitHub", "@EnzoTurpin"),
]
iy = cy + Inches(2.8)
for label, val in infos:
    add_text(slide, f"{label}: ", cx + Inches(0.25), iy, Inches(1.0), Inches(0.38),
             font_size=Pt(11), bold=True, color=YELLOW)
    add_text(slide, val, cx + Inches(1.25), iy, card_w - Inches(1.4), Inches(0.38),
             font_size=Pt(11), color=FG)
    iy += Inches(0.4)

# Carte 2 — Elouan
cx2 = Inches(6.95)
add_rect(slide, cx2, cy, card_w, card_h, fill_color=DARK2, line_color=YELLOW, line_width=Pt(2))
add_rect(slide, cx2, cy, card_w, Inches(0.05), fill_color=YELLOW)

add_rect(slide, cx2 + Inches(1.9), cy + Inches(0.25), Inches(1.7), Inches(1.7),
         fill_color=BG, line_color=YELLOW, line_width=Pt(2))
add_text(slide, "EC", cx2 + Inches(1.9), cy + Inches(0.35), Inches(1.7), Inches(1.5),
         font_size=Pt(44), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_text(slide, "ELOUAN CHE", cx2, cy + Inches(2.15), card_w, Inches(0.55),
         font_size=Pt(22), bold=True, color=FG, align=PP_ALIGN.CENTER)

infos2 = [
    ("Rôle", "Developer — Full-Stack"),
    ("Stack", "NestJS · React · PostgreSQL"),
    ("Email", "elouan.chedalleux@ynov.com"),
]
iy = cy + Inches(2.8)
for label, val in infos2:
    add_text(slide, f"{label}: ", cx2 + Inches(0.25), iy, Inches(1.0), Inches(0.38),
             font_size=Pt(11), bold=True, color=YELLOW)
    add_text(slide, val, cx2 + Inches(1.25), iy, card_w - Inches(1.4), Inches(0.38),
             font_size=Pt(11), color=FG)
    iy += Inches(0.4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DÉMO DU SITE (placeholder vidéo)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "DÉMO — APPLICATION WEB", "Slide 05")
add_footer(slide, 5)

# Zone vidéo placeholder
add_rect(slide, Inches(1.5), Inches(1.15), Inches(10.33), Inches(5.5),
         fill_color=DARK2, line_color=CYAN, line_width=Pt(2))

# Icône play
add_rect(slide, Inches(6.0), Inches(3.0), Inches(1.33), Inches(1.33),
         fill_color=BG, line_color=CYAN, line_width=Pt(3))
add_text(slide, "▶", Inches(6.0), Inches(2.9), Inches(1.33), Inches(1.4),
         font_size=Pt(40), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

add_text(slide, "INSÉRER LA VIDÉO DE DÉMO ICI",
         Inches(1.5), Inches(4.5), Inches(10.33), Inches(0.6),
         font_size=Pt(16), color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide,
         "Démo : Connexion · Création de board · Gestion des tâches · Partage · Favoris",
         Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.5),
         font_size=Pt(12), color=YELLOW, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DÉMO CLAUDE MCP (placeholder vidéo)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "DÉMO — CLAUDE MCP", "Slide 06")
add_footer(slide, 6)

add_rect(slide, Inches(1.5), Inches(1.15), Inches(10.33), Inches(5.5),
         fill_color=DARK2, line_color=YELLOW, line_width=Pt(2))

add_rect(slide, Inches(6.0), Inches(3.0), Inches(1.33), Inches(1.33),
         fill_color=BG, line_color=YELLOW, line_width=Pt(3))
add_text(slide, "▶", Inches(6.0), Inches(2.9), Inches(1.33), Inches(1.4),
         font_size=Pt(40), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_text(slide, "INSÉRER LA VIDÉO DE DÉMO MCP ICI",
         Inches(1.5), Inches(4.5), Inches(10.33), Inches(0.6),
         font_size=Pt(16), color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide,
         "Démo : Claude interagit avec l'API via MCP — création / modification de tâches en langage naturel",
         Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.5),
         font_size=Pt(12), color=CYAN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STACKS TECHNIQUES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "STACK TECHNIQUE", "Slide 07")
add_footer(slide, 7)

stacks = [
    {
        "cat": "FRONTEND",
        "color": CYAN,
        "items": [
            ("React 19", "UI — Concurrent features"),
            ("Vite 6", "Build tool ultra-rapide"),
            ("TypeScript 5", "Typage statique"),
            ("TailwindCSS v4", "Utility-first styling"),
            ("shadcn/ui", "Composants accessibles"),
            ("React Router v7", "Navigation SPA"),
            ("Zod", "Validation des schémas"),
        ]
    },
    {
        "cat": "BACKEND",
        "color": YELLOW,
        "items": [
            ("NestJS 11", "Framework Node.js DI"),
            ("Passport.js", "Authentification"),
            ("JWT dual-token", "Access 15min + Refresh 7j"),
            ("Google OAuth 2.0", "Connexion sociale"),
            ("Bcrypt", "Hashage des mots de passe"),
        ]
    },
    {
        "cat": "BASE DE DONNÉES",
        "color": RGBColor(0x00, 0xE8, 0x70),
        "items": [
            ("PostgreSQL ≥14", "Base relationnelle"),
            ("Prisma ORM v7", "Type-safe queries"),
            ("PG Adapter", "Connection pooling"),
        ]
    },
    {
        "cat": "OUTILS & DEVOPS",
        "color": RGBColor(0xA0, 0x60, 0xFF),
        "items": [
            ("Turbo v2", "Monorepo build system"),
            ("npm workspaces", "Gestion des packages"),
            ("ESLint 9 + Prettier", "Qualité de code"),
            ("Jest + Supertest", "Tests unitaires / API"),
        ]
    },
]

col_w = Inches(3.05)
for i, st in enumerate(stacks):
    cx = Inches(0.5) + i * (col_w + Inches(0.12))
    cy = Inches(1.1)
    col_h = Inches(5.6)
    add_rect(slide, cx, cy, col_w, col_h,
             fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
    add_rect(slide, cx, cy, col_w, Inches(0.05), fill_color=st["color"])

    add_text(slide, st["cat"], cx + Inches(0.1), cy + Inches(0.1),
             col_w - Inches(0.2), Inches(0.4),
             font_size=Pt(13), bold=True, color=st["color"])
    divider_line(slide, cy + Inches(0.58))

    iy = cy + Inches(0.68)
    for (name, desc) in st["items"]:
        add_text(slide, name, cx + Inches(0.15), iy,
                 col_w - Inches(0.25), Inches(0.3),
                 font_size=Pt(11.5), bold=True, color=FG)
        add_text(slide, desc, cx + Inches(0.15), iy + Inches(0.28),
                 col_w - Inches(0.25), Inches(0.3),
                 font_size=Pt(10), color=MUTED)
        iy += Inches(0.72)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "FONCTIONNALITÉS", "Slide 08")
add_footer(slide, 8)

features = [
    ("AUTH", "Authentification sécurisée",
     "Login email/password + Google OAuth\nJWT dual-token (access 15min / refresh 7j)\nAuto-refresh transparent"),
    ("BOARD", "Gestion de boards Kanban",
     "Création de boards & statuts colorés\nOrganisation des colonnes par drag\nBoards favoris persistants"),
    ("TASK", "Gestion des tâches",
     "Titre, description, deadline, priorité\nAssignation aux membres\nTracking d'activité / historique"),
    ("COLLAB", "Collaboration en équipe",
     "Partage par lien public (token)\nInvitations par email\nRôles : owner / admin / member"),
    ("MCP", "Intégration Claude MCP",
     "Claude accède à l'API via MCP\nCréation/modif tâches en langage naturel\nOutils MCP exposés et sécurisés"),
    ("UX", "Design eDEX-UI Terminal",
     "Charte graphique futuriste & cohérente\nMonospace · Scanlines · Glow effects\nFonts Share Tech Mono"),
]

col_w = Inches(4.0)
col_h = Inches(2.15)
for i, (tag, title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    cx = Inches(0.5) + col * (col_w + Inches(0.22))
    cy = Inches(1.1) + row * (col_h + Inches(0.18))

    add_rect(slide, cx, cy, col_w, col_h, fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
    add_rect(slide, cx, cy, Inches(0.07), col_h, fill_color=CYAN)

    add_rect(slide, cx + Inches(0.17), cy + Inches(0.1), Inches(0.7), Inches(0.32),
             fill_color=BG, line_color=CYAN, line_width=Pt(1))
    add_text(slide, tag, cx + Inches(0.17), cy + Inches(0.1), Inches(0.7), Inches(0.32),
             font_size=Pt(9), bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    add_text(slide, title, cx + Inches(0.95), cy + Inches(0.1),
             col_w - Inches(1.05), Inches(0.4),
             font_size=Pt(13), bold=True, color=FG)
    add_text(slide, desc, cx + Inches(0.17), cy + Inches(0.58),
             col_w - Inches(0.3), col_h - Inches(0.7),
             font_size=Pt(11), color=FG)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_corner_brackets(slide)
add_header_bar(slide, "CONCLUSION", "Slide 09")
add_footer(slide, 9)

# Résumé
add_rect(slide, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.6),
         fill_color=DARK2, line_color=MUTED, line_width=Pt(1))
add_rect(slide, Inches(0.5), Inches(1.1), Inches(7.5), Inches(0.05), fill_color=CYAN)

add_section_label(slide, "BILAN DU PROJET", Inches(0.65), Inches(1.2))

conclusions = [
    "Application Kanban full-stack fonctionnelle\net déployable en production",
    "Architecture robuste — monorepo Turbo,\nNestJS + React 19 + PostgreSQL",
    "Sécurité renforcée — JWT dual-token,\nrôles RBAC, validation Zod",
    "Intégration Claude MCP réussie —\ncréation de tâches en langage naturel",
    "Collaboration efficace en équipe de 2",
]
y = Inches(1.7)
for item in conclusions:
    add_rect(slide, Inches(0.7), y + Inches(0.18), Inches(0.1), Inches(0.1), fill_color=CYAN)
    add_text(slide, item, Inches(0.95), y, Inches(6.8), Inches(0.65),
             font_size=Pt(12.5), color=FG)
    y += Inches(0.88)

# Bloc droite — perspectives
add_rect(slide, Inches(8.25), Inches(1.1), Inches(4.5), Inches(2.7),
         fill_color=DARK2, line_color=YELLOW, line_width=Pt(1))
add_rect(slide, Inches(8.25), Inches(1.1), Inches(4.5), Inches(0.05), fill_color=YELLOW)
add_section_label(slide, "PERSPECTIVES", Inches(8.4), Inches(1.2))
persp = [
    "Déploiement sur cloud (Fly.io / Vercel)",
    "Notifications temps-réel (WebSocket)",
    "Tests E2E — Playwright",
    "Application mobile React Native",
]
py = Inches(1.7)
for p in persp:
    add_rect(slide, Inches(8.4), py + Inches(0.12), Inches(0.08), Inches(0.08), fill_color=YELLOW)
    add_text(slide, p, Inches(8.6), py, Inches(3.9), Inches(0.45),
             font_size=Pt(11.5), color=FG)
    py += Inches(0.52)

# Bloc QR / liens
add_rect(slide, Inches(8.25), Inches(4.05), Inches(4.5), Inches(2.65),
         fill_color=DARK2, line_color=CYAN, line_width=Pt(1))
add_rect(slide, Inches(8.25), Inches(4.05), Inches(4.5), Inches(0.05), fill_color=CYAN)
add_section_label(slide, "LIENS", Inches(8.4), Inches(4.15))
links = [
    ("GitHub", "github.com/EnzoTurpin/MCP"),
    ("Dev", "Enzo Turpin — @EnzoTurpin"),
    ("Dev", "Elouan Che — @ElouanChe"),
]
ly = Inches(4.65)
for label, val in links:
    add_text(slide, f"[{label}]", Inches(8.4), ly, Inches(0.9), Inches(0.38),
             font_size=Pt(11), bold=True, color=CYAN)
    add_text(slide, val, Inches(9.35), ly, Inches(3.2), Inches(0.38),
             font_size=Pt(11), color=FG)
    ly += Inches(0.5)

# Message final
add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.9), SLIDE_W - Inches(1), Inches(0.0))
add_text(slide, "[ MERCI POUR VOTRE ATTENTION ]",
         Inches(0.5), SLIDE_H - Inches(0.95), SLIDE_W - Inches(1), Inches(0.45),
         font_size=Pt(18), bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ─── Sauvegarde ────────────────────────────────────────────────────────────────
output_path = r"C:\Users\enzot\Desktop\Projets Ynov\Troisième Année\MCP\Jarvim_MCP_Presentation.pptx"
prs.save(output_path)
print(f"OK - Presentation generee : {output_path}")
print(f"   Nombre de slides : {len(prs.slides)}")
